import io
import random
from typing import List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt

from .models import SlideDeck


def collect_template_images(template_io: io.BytesIO, limit: int = 8) -> List[bytes]:
    prs = Presentation(template_io)
    blobs = []
    # Look for images in existing slides (if .pptx) and masters/layouts
    for sl in prs.slides:
        for shp in sl.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blobs.append(shp.image.blob)
                except Exception:
                    pass
            # background images are trickier; skip for simplicity
    # Deduplicate and cap
    uniq = []
    seen = set()
    for b in blobs:
        if b and (len(b) not in seen):
            uniq.append(b)
            seen.add(len(b))
        if len(uniq) >= limit:
            break
    return uniq


def _choose_layout(prs: Presentation, layout_hint: str):
    # Best effort mapping; fall back to first non-title layout
    names = [getattr(l, 'name', f'layout-{i}') or f'layout-{i}' for i, l in enumerate(prs.slide_layouts)]
    # Prefer common names
    priority = [
        "Title and Content",
        "Title and Vertical Text",
        "Two Content",
        "Title Only",
        "Section Header",
        "Content with Caption",
    ]
    # 1) exact match
    for i, l in enumerate(prs.slide_layouts):
        if names[i].lower() == (layout_hint or "").lower():
            return l
    # 2) priority
    for want in priority:
        for i, l in enumerate(prs.slide_layouts):
            if want.lower() in (names[i] or "").lower():
                return l
    # 3) fallback
    return prs.slide_layouts[1 if len(prs.slide_layouts) > 1 else 0]


def build_presentation(template_io: io.BytesIO, deck: SlideDeck) -> bytes:
    template_io.seek(0)
    prs = Presentation(template_io)

    # Build slides
    reusable_images = []
    try:
        reusable_images = collect_template_images(template_io)
    except Exception:
        pass

    rng = random.Random(42)

    for s in deck.slides:
        layout = _choose_layout(prs, s.layout_hint or "")
        slide = prs.slides.add_slide(layout)

        # Title
        if slide.shapes.title is not None:
            slide.shapes.title.text = s.title
            try:
                # Keep titles readable
                for p in slide.shapes.title.text_frame.paragraphs:
                    if p.font.size and p.font.size.pt > 48:
                        p.font.size = Pt(40)
            except Exception:
                pass

        # Body placeholder (first placeholder with text_frame)
        body = None
        for ph in slide.placeholders:
            if getattr(ph, "has_text_frame", False) and ph != slide.shapes.title:
                body = ph
                break
        if body is not None:
            tf = body.text_frame
            tf.clear()
            if s.bullets:
                tf.text = s.bullets[0]
                for b in s.bullets[1:]:
                    p = tf.add_paragraph()
                    p.text = b
                    p.level = 0

        # Optional: reuse template images in picture placeholders (if any)
        pics = [ph for ph in slide.placeholders if ph.placeholder_format and ph.placeholder_format.type == 18]  # 18=PICTURE placeholder
        if pics and reusable_images:
            chosen = rng.choice(reusable_images)
            # Insert as picture near the bounds of the first picture placeholder
            try:
                ph = pics[0]
                left, top, width, height = ph.left, ph.top, ph.width, ph.height
                slide.shapes.add_picture(io.BytesIO(chosen), left, top, width=width, height=height)
            except Exception:
                pass

        # Speaker notes
        try:
            notes = slide.notes_slide
        except Exception:
            notes = slide.notes_slide = slide.notes_slide
        if slide.notes_slide:
            slide.notes_slide.notes_text_frame.text = s.notes or ""

    bio = io.BytesIO()
    prs.save(bio)
    return bio.getvalue()